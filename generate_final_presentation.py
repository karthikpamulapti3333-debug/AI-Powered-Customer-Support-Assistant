import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    
    # 16:9 Widescreen Layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank layout
    
    # Color Palette Constants
    COLOR_BG = RGBColor(248, 250, 252)        # Slate Light Background #F8FAFC
    COLOR_PRIMARY = RGBColor(15, 23, 42)      # Deep Navy #0F172A
    COLOR_ACCENT = RGBColor(14, 165, 233)     # Electric Blue #0EA5E9
    COLOR_CARD = RGBColor(255, 255, 255)      # Pure White Card #FFFFFF
    COLOR_TEXT_MAIN = RGBColor(30, 41, 59)    # Dark Slate #1E293B
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)# Muted Gray #64748B
    COLOR_BORDER = RGBColor(226, 232, 240)    # Light Gray Border #E2E8F0
    COLOR_DARK_BG = RGBColor(15, 32, 67)      # Dark Navy Hero #0F2043

    def set_slide_background(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="AI CUSTOMER SUPPORT & TICKET MANAGEMENT SYSTEM"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT
        p_cat.font.name = "Arial"

        # Header Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.font.name = "Arial"

    def add_card(slide, left, top, width, height, title, items, icon="🔹", bg_color=COLOR_CARD, border_color=COLOR_BORDER):
        # Card Background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)

        # Content Box
        content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = content_box.text_frame
        tf.word_wrap = True

        if title:
            p_t = tf.paragraphs[0]
            p_t.text = f"{icon} {title}"
            p_t.font.size = Pt(16)
            p_t.font.bold = True
            p_t.font.color.rgb = COLOR_PRIMARY
            p_t.font.name = "Arial"
            p_t.space_after = Pt(10)

        for idx, item in enumerate(items):
            p = tf.add_paragraph() if (title or idx > 0) else tf.paragraphs[0]
            p.text = f"•  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_TEXT_MAIN
            p.font.name = "Arial"
            p.space_after = Pt(6)

    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_DARK_BG)

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "AI-Powered Customer Support and Ticket Management System"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.font.name = "Arial"
    p1.space_after = Pt(10)

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Automated AI Assistance, RAG Knowledge Base Retrieval & Intelligent Support Escalation"
    p1_sub.font.size = Pt(16)
    p1_sub.font.color.rgb = COLOR_ACCENT
    p1_sub.font.name = "Arial"

    meta_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.8), Inches(11.333), Inches(2.8))
    meta_card.fill.solid()
    meta_card.fill.fore_color.rgb = RGBColor(26, 43, 80)
    meta_card.line.color.rgb = COLOR_ACCENT

    tf_meta = slide1.shapes.add_textbox(Inches(1.3), Inches(4.0), Inches(10.7), Inches(2.4)).text_frame
    tf_meta.word_wrap = True

    meta_items = [
        ("Student Name:", "Pamulapati Karthik"),
        ("Roll Number:", "22N81A05K2"),
        ("Department:", "Computer Science & Engineering (CSE)"),
        ("College Name:", "Sree Dattha Group of Institutions"),
        ("Project Guide:", "Dr. A. Ramesh Kumar (HOD & Professor)")
    ]

    for label, val in meta_items:
        p = tf_meta.add_paragraph()
        run_l = p.add_run()
        run_l.text = f"{label:<18} "
        run_l.font.bold = True
        run_l.font.size = Pt(14)
        run_l.font.color.rgb = COLOR_ACCENT
        
        run_v = p.add_run()
        run_v.text = val
        run_v.font.size = Pt(14)
        run_v.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(6)

    add_notes(slide1, "Good morning respected guide, panel members, and peers. I am Pamulapati Karthik from the CSE Department. Today, I am presenting my project titled 'AI-Powered Customer Support and Ticket Management System'.")

    # ==========================================
    # SLIDE 2: Introduction
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_BG)
    add_header(slide2, "Introduction to AI-Based Customer Support")

    add_card(slide2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), "Overview of AI Support", [
        "AI customer support transforms static helpdesks into active 24/7 conversational agents.",
        "Combines Natural Language Processing (NLP) with domain-specific knowledge bases.",
        "Delivers instant assistance to website visitors without human intervention."
    ], "🌐")

    add_card(slide2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0), "Importance of Automation", [
        "Eliminates repetitive manual query handling by support staff.",
        "Scales customer service capacity seamlessly during high-volume spikes.",
        "Ensures consistent, accurate responses across all customer touchpoints."
    ], "🤖")

    add_card(slide2, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), "Benefits of Intelligent Chat", [
        "Zero customer registration friction on public landing pages.",
        "Reduces average response time from hours to under 2 seconds.",
        "Automates ticket escalation when queries require human expertise."
    ], "⚡")

    add_notes(slide2, "Customer service is moving towards instant automation. AI chatbots allow businesses to handle routine customer questions 24/7 without delays or human agent fatigue.")

    # ==========================================
    # SLIDE 3: Problem Statement
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_BG)
    add_header(slide3, "Problem Statement: Helpdesk Bottlenecks")

    add_card(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4), "Slow Response Times", [
        "Customers wait 24-48 hours for replies to simple email inquiries.",
        "Peak traffic leads to massive ticket backlogs and customer churn."
    ], "⏳")

    add_card(slide3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4), "High Operational Costs", [
        "Maintaining round-the-clock human support shifts is cost-prohibitive.",
        "Scaling human teams linearly with user growth is economically unsustainable."
    ], "💸")

    add_card(slide3, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4), "Limited Availability & Registration Barriers", [
        "Legacy support hours leave weekend and overnight gaps.",
        "Forcing users to register before asking basic questions causes high bounce rates."
    ], "🚪")

    add_card(slide3, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.4), "Lack of Automation & Poor UX", [
        "Rule-based bots fail on natural phrasing without generating tickets.",
        "Fragmented tools cause poor customer experience and staff burnout."
    ], "⚠️")

    add_notes(slide3, "Traditional helpdesks suffer from long wait times, high costs, and customer frustration. Legacy rule bots fail easily, and forcing user registration causes high drop-off.")

    # ==========================================
    # SLIDE 4: Objectives
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_BG)
    add_header(slide4, "Project Objectives & Key Goals")

    objectives = [
        ("Automate Customer Support", "Deploy an intelligent virtual assistant capable of resolving repetitive product and technical queries instantly."),
        ("Improve Response Time", "Cut average response time from hours to under 2 seconds using RAG Knowledge Base and real LLM integration."),
        ("Increase Customer Satisfaction", "Provide instant, accurate, natural language answers with markdown formatting and conversation memory."),
        ("Reduce Support Workload", "Automate up to 80% of routine inquiries, allowing human agents to focus on complex escalation cases."),
        ("Provide Continuous Support", "Ensure 24/7 availability with automated ticket creation whenever AI confidence falls below threshold.")
    ]

    for idx, (title, desc) in enumerate(objectives):
        top_pos = Inches(1.8 + idx * 1.0)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BORDER
        
        tf = slide4.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.1), Inches(11.3), Inches(0.65)).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"🎯 {idx+1}. {title}: "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_ACCENT
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_TEXT_MAIN

    add_notes(slide4, "Our project objectives focus on automating routine support, drastically improving response times, reducing support staff workload, and providing continuous 24/7 assistance.")

    # ==========================================
    # SLIDE 5: Existing System
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_BG)
    add_header(slide5, "Existing Support System vs. Limitations")

    add_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), "Existing Traditional Support Model", [
        "Relies heavily on manual ticketing portals and email forms.",
        "Mandatory user registration required prior to submitting queries.",
        "Fixed working hours resulting in weekend and overnight gaps.",
        "Static decision-tree rule bots with limited pattern matching.",
        "No natural language comprehension or context memory."
    ], "❌")

    add_card(slide5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), "Operational Consequences", [
        "High customer churn due to long wait times (24-48 hours).",
        "Higher support desk operational budget and agent burnout.",
        "High visitor bounce rates caused by mandatory login screens.",
        "Frequent frustration when rule bots loop on unrecognized phrases.",
        "Inconsistent ticket tracking and fragmented communication."
    ], "📉")

    add_notes(slide5, "Legacy systems suffer from static rule-based decision trees and manual ticket processing, creating high costs, slow resolution times, and poor user satisfaction.")

    # ==========================================
    # SLIDE 6: Proposed System
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_BG)
    add_header(slide6, "Proposed System Architecture & Solution")

    add_card(slide6, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), "AI-Powered Chatbot", [
        "Immediate guest access without customer login.",
        "Real LLM integration (OpenAI, Gemini, Ollama).",
        "Pre-LLM Knowledge Base RAG lookup.",
        "Session conversation memory & markdown rendering."
    ], "🤖")

    add_card(slide6, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0), "Automatic Ticket Generation", [
        "Confidence & sentiment threshold evaluation.",
        "Instant modal prompt when AI confidence is < 0.6.",
        "Generates unique Ticket ID (e.g. TICK-8F92A1B3).",
        "Stores structured tickets directly in DB."
    ], "🎫")

    add_card(slide6, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), "Admin Dashboard & Analytics", [
        "Single-auth Admin login (/admin/login).",
        "Master ticket queue management & replies.",
        "Knowledge Base FAQ CRUD management.",
        "Chart.js analytics & CSV report exporter."
    ], "⚙️")

    add_notes(slide6, "The proposed system combines an unauthenticated AI chatbot, automatic ticket generation for low-confidence queries, Knowledge Base integration, and a centralized Admin Dashboard.")

    # ==========================================
    # SLIDE 7: System Architecture Diagram
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLOR_BG)
    add_header(slide7, "System Architecture & High-Level Design")

    boxes = [
        (Inches(0.8), Inches(2.2), Inches(2.2), Inches(3.8), "Customer / Guest\nUser", "• User Interface\n• Instant Chat\n• Ticket Modal", COLOR_ACCENT),
        (Inches(3.3), Inches(2.2), Inches(2.4), Inches(3.8), "Flask Application\n(Backend Engine)", "• App Factory\n• Routing Controllers\n• Jinja2 Templates", COLOR_PRIMARY),
        (Inches(6.0), Inches(2.2), Inches(2.4), Inches(3.8), "AI Engine & RAG\nModule", "• KB RAG Search\n• LLM Client\n• Intent Detector", COLOR_DARK_BG),
        (Inches(8.7), Inches(2.2), Inches(2.0), Inches(3.8), "Database Layer\n(SQLAlchemy)", "• SQLite (Dev)\n• PostgreSQL (Prod)\n• Auto Schemas", COLOR_PRIMARY),
        (Inches(11.0), Inches(2.2), Inches(1.5), Inches(3.8), "Admin\nDashboard", "• Flask-Login\n• Ticket Queue\n• CSV Export", COLOR_ACCENT)
    ]

    for left, top, width, height, title, body, color in boxes:
        shape = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(2.0)

        tf = slide7.shapes.add_textbox(left + Inches(0.1), top + Inches(0.2), width - Inches(0.2), height - Inches(0.4)).text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)

        p_body = tf.add_paragraph()
        p_body.text = body
        p_body.font.size = Pt(11)
        p_body.font.color.rgb = COLOR_TEXT_MAIN
        p_body.alignment = PP_ALIGN.LEFT

    add_notes(slide7, "The system architecture connects Customer UI, Flask Application Backend, AI Engine with RAG search, Database layer, Ticket Management System, and Administrator Dashboard.")

    # ==========================================
    # SLIDE 8: System Modules
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLOR_BG)
    add_header(slide8, "System Modules Breakdown")

    modules = [
        ("Customer Module", ["Instant guest access on homepage.", "Zero login or registration required.", "Interactive Bootstrap 5 dark theme."], "👤"),
        ("Chat Module", ["Real-time streaming typing animation.", "Marked.js markdown rendering for code.", "Session conversation memory context."], "💬"),
        ("Ticket Management Module", ["Automated Ticket ID generation (TICK-XXXXXXXX).", "Captures Name, Email, Phone, Priority.", "Database persistence & status tracking."], "🎫"),
        ("Knowledge Base Module", ["Pre-LLM RAG FAQ pattern lookup.", "Dynamic CRUD management by Admin.", "Reduces external LLM API costs."], "📚"),
        ("Administrator Module", ["Flask-Login protected authentication.", "Ticket reply threading & status controls.", "Chart.js analytics & CSV report exporter."], "⚙️")
    ]

    for idx, (title, items, icon) in enumerate(modules):
        top_pos = Inches(1.8 + idx * 1.0)
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BORDER
        
        tf = slide8.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.1), Inches(11.3), Inches(0.65)).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{icon} {title}: "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_PRIMARY
        
        r2 = p.add_run()
        r2.text = " | ".join(items)
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_TEXT_MAIN

    add_notes(slide8, "The system consists of five modular components: Customer Module, Chat Module, Ticket Management Module, Knowledge Base Module, and Administrator Module.")

    # ==========================================
    # SLIDE 9: Database Design
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, COLOR_BG)
    add_header(slide9, "Database Design & Relational Tables")

    tables = [
        ("Admin Table", ["id (PK)", "email", "username", "password_hash", "role", "created_at"]),
        ("ChatSessions Table", ["id (PK)", "session_id (Index)", "title", "status", "created_at"]),
        ("Messages Table", ["id (PK)", "session_id (FK)", "sender", "content", "intent", "sentiment"]),
        ("Tickets Table", ["id (PK)", "ticket_code (Unique)", "customer_name", "email", "phone", "status"]),
        ("KnowledgeBase Table", ["id (PK)", "question", "answer", "category", "is_published"]),
        ("ActivityLogs Table", ["id (PK)", "admin_id (FK)", "action", "details", "timestamp"])
    ]

    for idx, (tname, cols) in enumerate(tables):
        col = idx % 3
        row = idx // 3
        left = Inches(0.8 + col * 3.9)
        top = Inches(1.8 + row * 2.6)
        add_card(slide9, left, top, Inches(3.7), Inches(2.4), tname, cols, "🗄️")

    add_notes(slide9, "The database design features 6 core tables: Admin, ChatSessions, Messages, Tickets, KnowledgeBase, and ActivityLogs, with automatic schema creation via db.create_all().")

    # ==========================================
    # SLIDE 10: Workflow
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, COLOR_BG)
    add_header(slide10, "End-to-End System Workflow")

    steps = [
        ("Step 1: Customer Submits Query", "Customer opens website and types query into instant ChatGPT interface."),
        ("Step 2: AI Processes Request", "System performs pre-LLM Knowledge Base RAG search & intent classification."),
        ("Step 3: AI Generates Response", "Renders answer in markdown format with session conversation memory."),
        ("Step 4: Ticket Generation", "If AI confidence is low (< 0.6), system prompts guest & generates Ticket ID."),
        ("Step 5: Administrator Resolves", "Administrator reviews ticket in console, posts reply, and updates status.")
    ]

    for idx, (title, desc) in enumerate(steps):
        top_pos = Inches(1.8 + idx * 1.0)
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_ACCENT if idx == 3 else COLOR_BORDER
        card.line.width = Pt(1.5)

        tf = slide10.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.1), Inches(11.3), Inches(0.65)).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{title}: "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_PRIMARY
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_TEXT_MAIN

    add_notes(slide10, "The 5-step workflow connects customer submission, AI processing, AI response, automated ticket generation when necessary, and administrator resolution.")

    # ==========================================
    # SLIDE 11: Core System Features
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, COLOR_BG)
    add_header(slide11, "Core System Features & Capabilities")

    add_card(slide11, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4), "Interactive Chat Interface", [
        "Markdown formatting for code blocks and tables.",
        "Streaming typing animation & auto-scroll."
    ], "💬")

    add_card(slide11, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4), "Automated Ticket Creation", [
        "Unique Ticket ID generation (TICK-XXXXXXXX).",
        "Captures contact info without user accounts."
    ], "🎫")

    add_card(slide11, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4), "Admin Dashboard & Search", [
        "Real-time ticket queue status controls.",
        "Search & filter tickets by status or priority."
    ], "📊")

    add_card(slide11, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.4), "Analytics & Data Export", [
        "Chart.js ticket status doughnut graph.",
        "One-click CSV report export (/admin/export/tickets)."
    ], "📈")

    add_notes(slide11, "Core features include real-time chat interface, automated ticket creation, administrative search and dashboard, Chart.js analytics, and CSV report export.")

    # ==========================================
    # SLIDE 12: Tools and Technologies Used
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, COLOR_BG)
    add_header(slide12, "Tools & Technology Stack")

    stacks = [
        ("Backend & ORM", ["Python 3.12", "Flask Framework", "Flask-SQLAlchemy", "Flask-Login", "Gunicorn WSGI"], "🐍"),
        ("Frontend Technologies", ["HTML5 & CSS3", "Bootstrap 5", "Vanilla JavaScript", "Marked.js Markdown", "Jinja2 Templating"], "🎨"),
        ("Databases", ["SQLite (Development)", "PostgreSQL (Production)", "Zero-Migration Setup", "db.create_all() Auto Schema"], "💾"),
        ("Deployment & AI", ["OpenAI-compatible API", "Gemini / Ollama", "GitHub Repository", "Render Cloud Platform"], "☁️")
    ]

    for idx, (title, items, icon) in enumerate(stacks):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.8 + row * 2.6)
        add_card(slide12, left, top, Inches(5.6), Inches(2.4), title, items, icon)

    add_notes(slide12, "The tech stack relies on Python 3.12, Flask, SQLAlchemy, SQLite/PostgreSQL, Bootstrap 5, Jinja2, GitHub, Render, and OpenAI-compatible APIs.")

    # ==========================================
    # SLIDE 13: Challenges Faced & Solutions
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, COLOR_BG)
    add_header(slide13, "Technical Challenges & Solutions")

    add_card(slide13, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4), "Render Deployment & AppImportError", [
        "Issue: Gunicorn could not locate Flask app instance.",
        "Solution: Explicitly exposed app = create_app() in app/__init__.py."
    ], "🛠️")

    add_card(slide13, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4), "Database Migration & Schema Mismatch", [
        "Issue: Flask-Migrate revision mismatches during cloud deploy.",
        "Solution: Replaced Alembic with zero-migration db.create_all()."
    ], "💾")

    add_card(slide13, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4), "API Integration & Local Fallbacks", [
        "Issue: Intermittent external LLM rate limits.",
        "Solution: Built pre-LLM RAG KB search and intelligent local simulator."
    ], "🔌")

    add_card(slide13, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.4), "UI UX & Performance Optimization", [
        "Issue: Streaming chat rendering and mobile responsiveness.",
        "Solution: Implemented Bootstrap 5 glassmorphism and Marked.js."
    ], "⚡")

    add_notes(slide13, "Challenges included resolving Gunicorn AppImportErrors on Render, removing database migration conflicts, building local LLM fallback clients, and optimizing UI performance.")

    # ==========================================
    # SLIDE 14: Results and Outcomes
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, COLOR_BG)
    add_header(slide14, "Results & Project Outcomes")

    add_card(slide14, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), "Reduced Response Time", [
        "Queries resolved in < 2 seconds.",
        "Eliminated 24-48 hour email delays.",
        "Instant guest access without login."
    ], "⚡")

    add_card(slide14, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0), "Improved Operational Efficiency", [
        "Automates 80% of routine questions.",
        "Seamless ticket escalation for low confidence.",
        "Saves human agent workload."
    ], "📈")

    add_card(slide14, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), "Simplified Management Process", [
        "Centralized Admin Dashboard.",
        "Real-time Chart.js analytics.",
        "One-click CSV report exporter."
    ], "🎯")

    add_notes(slide14, "Key outcomes include reducing response time to under 2 seconds, automating routine queries, improving customer experience, and simplifying helpdesk management.")

    # ==========================================
    # SLIDE 15: Future Enhancements
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, COLOR_BG)
    add_header(slide15, "Future Enhancements & Roadmap")

    add_card(slide15, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4), "Voice Support & Audio Input", [
        "Integrate Speech-to-Text for real-time voice interactions.",
        "Support audio message processing."
    ], "🎙️")

    add_card(slide15, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4), "Multiple Language Support", [
        "Implement automatic language translation for global users.",
        "Localized UI themes and responses."
    ], "🌐")

    add_card(slide15, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4), "Mobile Applications", [
        "Develop native Android and iOS mobile applications.",
        "Push notifications for ticket updates."
    ], "📱")

    add_card(slide15, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.4), "Advanced Analytics & AI Capabilities", [
        "Predictive churn analysis and SLA breach warnings.",
        "Enhanced fine-tuned domain LLMs."
    ], "🔮")

    add_notes(slide15, "Future enhancements include adding voice interaction, multi-language support, native mobile apps, and advanced predictive analytics.")

    # ==========================================
    # SLIDE 16: Conclusion
    # ==========================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16, COLOR_BG)
    add_header(slide16, "Conclusion & Summary")

    add_card(slide16, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Project Summary", [
        "Successfully developed and deployed a modern, production-ready AI Customer Support & Ticket Management Platform.",
        "Demonstrated zero-friction guest AI assistance with RAG Knowledge Base search and real LLM integration.",
        "Implemented seamless automated ticket escalation for low-confidence queries without requiring user accounts.",
        "Delivered a secure, single-auth Admin Management Console powered by Flask-Login with live Chart.js analytics and CSV report export.",
        "Achieved 100% clean automated unit test execution and zero-migration cloud deployment on Render via Gunicorn."
    ], "🎯")

    add_notes(slide16, "In conclusion, this project demonstrates how combining Large Language Models, Knowledge Base RAG retrieval, and clean web engineering can transform customer support operations.")

    # ==========================================
    # SLIDE 17: Thank You Slide
    # ==========================================
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17, COLOR_DARK_BG)

    ty_box = slide17.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(1.8))
    tf_ty = ty_box.text_frame
    tf_ty.word_wrap = True
    
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "Thank You!"
    p_ty.font.size = Pt(44)
    p_ty.font.bold = True
    p_ty.font.color.rgb = RGBColor(255, 255, 255)
    p_ty.alignment = PP_ALIGN.CENTER
    
    p_sub = tf_ty.add_paragraph()
    p_sub.text = "Questions & Answers / Project Review Session"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_ACCENT
    p_sub.alignment = PP_ALIGN.CENTER

    contact_card = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(3.8), Inches(8.333), Inches(2.6))
    contact_card.fill.solid()
    contact_card.fill.fore_color.rgb = RGBColor(26, 43, 80)
    contact_card.line.color.rgb = COLOR_ACCENT

    tf_c = slide17.shapes.add_textbox(Inches(2.7), Inches(4.0), Inches(7.933), Inches(2.2)).text_frame
    tf_c.word_wrap = True

    c_info = [
        ("Presenter:", "Pamulapati Karthik (Roll No: 22N81A05K2)"),
        ("Department:", "Computer Science & Engineering (CSE)"),
        ("College:", "Sree Dattha Group of Institutions"),
        ("GitHub Repository:", "github.com/karthikpamulapti3333-debug/AI-Powered-Customer-Support-Assistant")
    ]

    for label, val in c_info:
        p = tf_c.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{label:<18} "
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_ACCENT
        
        r2 = p.add_run()
        r2.text = val
        r2.font.size = Pt(13)
        r2.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

    add_notes(slide17, "Thank you very much for your time, guidance, and attention throughout this presentation. I am now open to your questions and project review.")

    # Save presentation
    output_filename = "AI_Customer_Support_Presentation_Final.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'")

    # Also copy to shared_sadre/Presentation/Project_Presentation.pptx if folder exists
    sadre_pptx = os.path.join(os.path.dirname(__file__), "shared_sadre", "Presentation", "Project_Presentation.pptx")
    if os.path.exists(os.path.dirname(sadre_pptx)):
        shutil.copy2(output_filename, sadre_pptx)
        print(f"Updated presentation in shared_sadre: {sadre_pptx}")

    return output_filename

if __name__ == "__main__":
    create_deck()
