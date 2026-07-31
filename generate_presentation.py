import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    
    # 16:9 Widescreen Layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank layout
    
    # Color Palette Constants
    COLOR_BG = RGBColor(248, 250, 252)        # Slate Light Background #F8FAFC
    COLOR_PRIMARY = RGBColor(15, 23, 42)      # Navy / Deep Slate #0F172A
    COLOR_ACCENT = RGBColor(14, 165, 233)     # Electric Blue #0EA5E9
    COLOR_CARD = RGBColor(255, 255, 255)      # Pure White Card #FFFFFF
    COLOR_TEXT_MAIN = RGBColor(30, 41, 59)    # Dark Slate #1E293B
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)# Muted Gray #64748B
    COLOR_BORDER = RGBColor(226, 232, 240)    # Light Gray Border #E2E8F0
    COLOR_DARK_BG = RGBColor(15, 32, 67)      # Dark Navy Hero #0F2043

    def set_slide_background(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="AI-POWERED CUSTOMER SUPPORT PLATFORM"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT
        p_cat.font.name = "Arial"

        # Header Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.font.name = "Arial"

    def add_card(slide, left, top, width, height, title, items, icon="🔹", bg_color=COLOR_CARD, border_color=COLOR_BORDER):
        # Card Background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)

        # Content Box
        content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = content_box.text_frame
        tf.word_wrap = True

        if title:
            p_t = tf.paragraphs[0]
            p_t.text = f"{icon} {title}"
            p_t.font.size = Pt(16)
            p_t.font.bold = True
            p_t.font.color.rgb = COLOR_PRIMARY
            p_t.font.name = "Arial"
            p_t.space_after = Pt(10)

        for idx, item in enumerate(items):
            p = tf.add_paragraph() if (title or idx > 0) else tf.paragraphs[0]
            p.text = f"•  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_TEXT_MAIN
            p.font.name = "Arial"
            p.space_after = Pt(6)

    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_DARK_BG)

    # Title Banner Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "AI-Powered Customer Support & Ticket Management System"
    p1.font.size = Pt(34)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.font.name = "Arial"
    p1.space_after = Pt(10)

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "A Modern, Production-Ready Customer Support Platform with Instant AI Assistance & Automated Escalation"
    p1_sub.font.size = Pt(16)
    p1_sub.font.color.rgb = COLOR_ACCENT
    p1_sub.font.name = "Arial"

    # Metadata Grid (Card Style)
    meta_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.8), Inches(11.333), Inches(2.8))
    meta_card.fill.solid()
    meta_card.fill.fore_color.rgb = RGBColor(26, 43, 80)
    meta_card.line.color.rgb = COLOR_ACCENT

    tf_meta = slide1.shapes.add_textbox(Inches(1.3), Inches(4.0), Inches(10.7), Inches(2.4)).text_frame
    tf_meta.word_wrap = True

    meta_items = [
        ("Student Name:", "Pamulapati Karthik"),
        ("Roll Number:", "22N81A05K2"),
        ("Department:", "Computer Science & Engineering (CSE)"),
        ("College Name:", "Sree Dattha Group of Institutions"),
        ("Project Guide:", "Dr. A. Ramesh Kumar (HOD & Professor)")
    ]

    for label, val in meta_items:
        p = tf_meta.add_paragraph()
        run_l = p.add_run()
        run_l.text = f"{label:<18} "
        run_l.font.bold = True
        run_l.font.size = Pt(14)
        run_l.font.color.rgb = COLOR_ACCENT
        
        run_v = p.add_run()
        run_v.text = val
        run_v.font.size = Pt(14)
        run_v.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(6)

    add_notes(slide1, "Good morning/afternoon respected guide, panel members, and peers. I am Pamulapati Karthik from the Department of Computer Science & Engineering. Today, I am excited to present my project titled 'AI-Powered Customer Support and Ticket Management System'. This project addresses the key challenges in modern customer service by combining instant AI chatbot assistance with automated ticket management.")

    # ==========================================
    # SLIDE 2: Introduction
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_BG)
    add_header(slide2, "Introduction to Modern AI Customer Support")

    add_card(slide2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), "Customer Support Overview", [
        "Customer support is a critical business function that directly drives retention.",
        "Traditional support relies heavily on phone calls, static emails, and manual queues.",
        "High volume leads to delayed responses and customer dissatisfaction."
    ], "🌐")

    add_card(slide2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0), "Role of Artificial Intelligence", [
        "AI enables instant natural language understanding and automated query resolution.",
        "Processes queries 24/7 without fatigue or queue delays.",
        "Seamlessly searches internal Knowledge Bases before calling LLMs."
    ], "🤖")

    add_card(slide2, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), "Need for Intelligent Automation", [
        "Eliminates initial friction by offering instant guest assistance without login barriers.",
        "Automates ticket escalation when AI confidence is insufficient.",
        "Drastically reduces operational support costs while improving satisfaction."
    ], "🚀")

    add_notes(slide2, "Customer support is the frontline of customer retention. However, traditional manual support models struggle with volume spikes, high cost, and slow turnaround times. AI technology changes this dynamic by offering instant, 24/7 contextual answers. By integrating Large Language Models and Retrieval-Augmented Generation, businesses can resolve up to 80% of repetitive inquiries automatically while escalating complex cases seamlessly.")

    # ==========================================
    # SLIDE 3: Problem Statement
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_BG)
    add_header(slide3, "Problem Statement: Challenges in Support Operations")

    add_card(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4), "High Operational Expenses", [
        "Maintaining 24/7 human agent teams requires substantial budget overhead.",
        "Scaling human teams linearly with user growth is economically unsustainable."
    ], "💸")

    add_card(slide3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4), "Delayed Response & Resolution Times", [
        "Customers wait hours or days for responses to routine inquiries.",
        "Peak hour surges cause massive ticket backlogs and customer churn."
    ], "⏳")

    add_card(slide3, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4), "High Customer Onboarding Friction", [
        "Forcing users to register/log in just to ask a simple pre-sale question leads to high drop-off.",
        "Inconsistent answers across human support staff degrade brand trust."
    ], "🚪")

    add_card(slide3, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.4), "Lack of Unified Ticket Escalation", [
        "Traditional chatbots fail silently without creating actionable tickets.",
        "Support managers lack central analytics to monitor resolution velocity."
    ], "⚠️")

    add_notes(slide3, "Businesses face a dilemma: customers expect instant 24/7 assistance, but scaling human support teams is prohibitively expensive. Furthermore, forcing customers to create an account just to ask a basic pre-sales question leads to high drop-off. When basic chatbots fail, they leave users stranded without auto-generating structured support tickets. Our project directly solves this problem.")

    # ==========================================
    # SLIDE 4: Objectives
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_BG)
    add_header(slide4, "Project Objectives & Key Goals")

    objectives = [
        ("Automate Customer Support", "Deploy an intelligent virtual assistant capable of resolving repetitive product and technical queries instantly."),
        ("Zero-Friction Access", "Enable immediate guest access on homepage load without requiring customer login or account creation."),
        ("Drastic Response Reduction", "Cut average query response time from hours to under 2 seconds using RAG Knowledge Base and real LLM integration."),
        ("Automated Ticket Generation", "Automatically prompt and generate structured support tickets with unique IDs when AI confidence is under threshold."),
        ("Centralized Admin Console", "Provide administrators with full ticket lifecycle control, Chart.js analytics, and CSV report export capability.")
    ]

    for idx, (title, desc) in enumerate(objectives):
        top_pos = Inches(1.8 + idx * 1.0)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BORDER
        
        tf = slide4.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.1), Inches(11.3), Inches(0.65)).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"🎯 {idx+1}. {title}: "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_ACCENT
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_TEXT_MAIN

    add_notes(slide4, "The primary objectives of this project are: first, automate responses to routine customer questions; second, eliminate login barriers for visitors; third, reduce response times from hours to seconds; fourth, guarantee that unresolved issues cleanly escalate to support tickets; and fifth, equip managers with a centralized dashboard and analytics.")

    # ==========================================
    # SLIDE 5: Existing System vs Limitations
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_BG)
    add_header(slide5, "Existing Support Systems vs. Key Limitations")

    add_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), "Existing Traditional Support Model", [
        "Relies heavily on manual ticketing portals and email forms.",
        "Mandatory user registration required prior to submitting queries.",
        "Fixed working hours resulting in weekend and overnight gaps.",
        "Static decision-tree rule bots with limited pattern matching.",
        "No natural language comprehension or context memory."
    ], "❌")

    add_card(slide5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), "Operational Consequences", [
        "High customer churn due to long wait times (24-48 hours).",
        "Higher support desk operational budget and agent burnout.",
        "High visitor bounce rates caused by mandatory login screens.",
        "Frequent frustration when rule bots loop on unrecognized phrases.",
        "Inconsistent ticket tracking and fragmented communication."
    ], "📉")

    add_notes(slide5, "Existing legacy systems rely on manual ticketing or basic decision-tree rule bots. Rule bots break whenever a user phrases a query naturally, while manual queues force users to wait up to 48 hours for simple answers. Additionally, requiring account registration before answering a basic question creates unnecessary friction for users.")

    # ==========================================
    # SLIDE 6: Proposed System
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_BG)
    add_header(slide6, "Proposed System Architecture & Solution")

    add_card(slide6, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), "Instant AI Chatbot Engine", [
        "Immediate guest access without customer login.",
        "Real LLM integration (OpenAI, Gemini, Ollama).",
        "Pre-LLM Knowledge Base RAG lookup.",
        "Session conversation memory & markdown rendering."
    ], "🤖")

    add_card(slide6, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0), "Automated Ticket Escalation", [
        "Confidence & sentiment threshold evaluation.",
        "Instant modal prompt when AI confidence is < 0.6.",
        "Generates unique Ticket ID (e.g. TICK-8F92A1B3).",
        "Stores structured tickets directly in DB."
    ], "🎫")

    add_card(slide6, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), "Admin Management Console", [
        "Single-auth Admin login (/admin/login).",
        "Master ticket lifecycle management & replies.",
        "Knowledge Base FAQ CRUD management.",
        "Chart.js analytics & CSV report exporter."
    ], "⚙️")

    add_notes(slide6, "Our proposed solution introduces a unified AI platform: 1) Instant guest access to a ChatGPT-style assistant powered by real LLM APIs and local Knowledge Base RAG; 2) Intelligent automated ticket escalation whenever the AI's confidence is under threshold; and 3) A secure Admin Management Console where managers review analytics, handle tickets, and update published FAQs.")

    # ==========================================
    # SLIDE 7: System Architecture Diagram
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLOR_BG)
    add_header(slide7, "System Architecture & High-Level Design")

    # Architecture Box Diagram Elements
    boxes = [
        (Inches(0.8), Inches(2.2), Inches(2.2), Inches(3.8), "Customer / Guest\nVisitor", "• Instant Chat\n• Ticket Modal\n• Markdown UI", COLOR_ACCENT),
        (Inches(3.3), Inches(2.2), Inches(2.4), Inches(3.8), "Flask Web App\n(Backend Engine)", "• App Factory\n• Routing Controllers\n• Jinja2 Templates", COLOR_PRIMARY),
        (Inches(6.0), Inches(2.2), Inches(2.4), Inches(3.8), "AI Engine & RAG\nModule", "• KB Search RAG\n• LLM Client\n• Intent Classification", COLOR_DARK_BG),
        (Inches(8.7), Inches(2.2), Inches(2.0), Inches(3.8), "Database Layer\n(SQLAlchemy)", "• SQLite (Dev)\n• PostgreSQL (Prod)\n• Auto Schema Creation", COLOR_PRIMARY),
        (Inches(11.0), Inches(2.2), Inches(1.5), Inches(3.8), "Admin\nConsole", "• Flask-Login\n• Analytics\n• CSV Export", COLOR_ACCENT)
    ]

    for left, top, width, height, title, body, color in boxes:
        shape = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(2.0)

        tf = slide7.shapes.add_textbox(left + Inches(0.1), top + Inches(0.2), width - Inches(0.2), height - Inches(0.4)).text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)

        p_body = tf.add_paragraph()
        p_body.text = body
        p_body.font.size = Pt(11)
        p_body.font.color.rgb = COLOR_TEXT_MAIN
        p_body.alignment = PP_ALIGN.LEFT

    add_notes(slide7, "Here we see the end-to-end system architecture. Public visitors interact with the Web App frontend. The Flask Backend routes messages to the AI Engine, which first performs a Retrieval-Augmented Generation search against local Knowledge Base FAQs before calling external LLMs. The database layer uses SQLAlchemy for SQLite or PostgreSQL, and the Admin Console provides secure session control via Flask-Login.")

    # ==========================================
    # SLIDE 8: Technology Stack
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLOR_BG)
    add_header(slide8, "Technology Stack & Frameworks")

    stacks = [
        ("Backend Framework", ["Python 3.12", "Flask App Factory", "Flask-SQLAlchemy", "Flask-Login", "Gunicorn WSGI"], "🐍"),
        ("Frontend & UI", ["HTML5 & CSS3", "Bootstrap 5", "Vanilla JavaScript", "Marked.js Markdown", "Glassmorphism UI"], "🎨"),
        ("Database & Storage", ["SQLite (Development)", "PostgreSQL (Production)", "db.create_all() Auto Schema", "Zero-Migration Setup"], "💾"),
        ("AI & Infrastructure", ["OpenAI GPT API", "Google Gemini API", "Ollama Local Model", "Render Cloud Deployment"], "☁️")
    ]

    for idx, (title, items, icon) in enumerate(stacks):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.8 + row * 2.6)
        add_card(slide8, left, top, Inches(5.6), Inches(2.4), title, items, icon)

    add_notes(slide8, "Our technical stack is modern and lightweight: Python 3.12 and Flask for the backend API; Flask-SQLAlchemy for database ORM with zero-migration auto-schema creation; Bootstrap 5 and Glassmorphism dark mode for the frontend UI; and multi-provider AI support covering OpenAI, Gemini, Ollama, and local fallbacks, deployed on Render with Gunicorn.")

    # ==========================================
    # SLIDE 9: Database Schema Design
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, COLOR_BG)
    add_header(slide9, "Database Schema & Entity Model")

    tables = [
        ("Admin", ["id (PK)", "email", "username", "password_hash", "role", "created_at"]),
        ("ChatSession", ["id (PK)", "session_id (Index)", "title", "status", "created_at", "updated_at"]),
        ("Message", ["id (PK)", "session_id (FK)", "sender (USER/AI)", "content", "intent", "sentiment"]),
        ("Ticket", ["id (PK)", "ticket_code (Unique)", "customer_name", "email", "phone", "subject", "status"]),
        ("KnowledgeBase", ["id (PK)", "question", "answer", "category", "is_published", "created_at"]),
        ("ActivityLog", ["id (PK)", "admin_id (FK)", "action", "details", "ip_address", "timestamp"])
    ]

    for idx, (tname, cols) in enumerate(tables):
        col = idx % 3
        row = idx // 3
        left = Inches(0.8 + col * 3.9)
        top = Inches(1.8 + row * 2.6)
        add_card(slide9, left, top, Inches(3.7), Inches(2.4), f"Table: {tname}", cols, "🗄️")

    add_notes(slide9, "The database relational model comprises 6 core tables: Admin (stores hashed admin credentials), ChatSession and Message (track conversation context), Ticket (stores guest tickets with generated Ticket IDs), KnowledgeBase (stores published FAQs for RAG search), and ActivityLog (audit trail for admin actions). Schema creation is completely automated on app startup via db.create_all().")

    # ==========================================
    # SLIDE 10: System Workflow
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, COLOR_BG)
    add_header(slide10, "End-to-End Execution Workflow")

    steps = [
        ("Step 1: Guest Inquiry", "Customer opens website and types query into instant ChatGPT interface."),
        ("Step 2: Pre-LLM RAG Search", "System searches local KnowledgeBase FAQs for exact pattern match."),
        ("Step 3: AI LLM Generation", "If no FAQ match, queries configurable LLM (OpenAI/Gemini) with session history."),
        ("Step 4: Ticket Escalation", "If confidence < 0.6, system prompts guest to open ticket & generates Ticket ID."),
        ("Step 5: Admin Resolution", "Administrator reviews ticket in console, posts reply, and updates status.")
    ]

    for idx, (title, desc) in enumerate(steps):
        top_pos = Inches(1.8 + idx * 1.0)
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_ACCENT if idx == 3 else COLOR_BORDER
        card.line.width = Pt(1.5)

        tf = slide10.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.1), Inches(11.3), Inches(0.65)).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{title}: "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_PRIMARY
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_TEXT_MAIN

    add_notes(slide10, "This slide outlines the 5-step operational workflow. 1) The visitor asks a question; 2) The RAG module checks Knowledge Base FAQs; 3) If unmatched, the LLM processes the query with session context; 4) If AI confidence is low, a ticket modal opens automatically; 5) The administrator manages and resolves the ticket in the admin console.")

    # ==========================================
    # SLIDE 11: Customer Module
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, COLOR_BG)
    add_header(slide11, "Customer Experience & Public Module")

    add_card(slide11, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), "Instant ChatGPT Interface", [
        "Embedded directly on homepage.",
        "Zero registration or login required.",
        "Real-time streaming typing animation.",
        "Marked.js markdown rendering.",
        "Auto-scroll & clear conversation."
    ], "💬")

    add_card(slide11, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0), "Guest Ticket Submission", [
        "Triggered automatically or manually.",
        "Captures Name, Email, Phone, Subject, Description, Priority.",
        "Generates unique Ticket ID (TICK-XXXXXXXX).",
        "Saves directly to DB without user account."
    ], "🎫")

    add_card(slide11, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), "Interactive UX Design", [
        "Sleek Glassmorphic dark theme.",
        "Pre-filled quick question prompts.",
        "Copy response button & code highlight.",
        "Fully responsive on mobile & desktop."
    ], "✨")

    add_notes(slide11, "The Customer Module is designed for zero friction: public visitors immediately interact with an intuitive ChatGPT-style chatbot on page load. When an issue requires human attention, a simple ticket form modal captures contact details and issues a unique Ticket ID without asking the user to create an account.")

    # ==========================================
    # SLIDE 12: Admin Module
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, COLOR_BG)
    add_header(slide12, "Admin Console & Operations Module")

    add_card(slide12, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), "Authentication & Dashboard", [
        "Flask-Login protected access (/admin/login).",
        "Seeded credentials (admin@example.com).",
        "Stat cards: Total, Open, Closed Tickets, AI Conversations, Messages.",
        "Recent ticket queue overview."
    ], "⚙️")

    add_card(slide12, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0), "Ticket & FAQ Management", [
        "Search, filter by status, view details.",
        "Post admin replies & update ticket status.",
        "Delete resolved or spam tickets.",
        "Publish, edit, and delete Knowledge Base FAQs."
    ], "📋")

    add_card(slide12, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), "Analytics & Data Export", [
        "Chart.js ticket status doughnut graph.",
        "Category distribution bar chart.",
        "7-day chat activity timeline.",
        "One-click CSV report exporter (/admin/export/tickets)."
    ], "📊")

    add_notes(slide12, "The Admin Module provides complete operational control: secure Flask-Login authentication, a comprehensive dashboard displaying real-time metrics, ticket threading and status updates, Knowledge Base FAQ publishing, Chart.js analytics graphs, and one-click CSV report exports.")

    # ==========================================
    # SLIDE 13: Core System Features
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, COLOR_BG)
    add_header(slide13, "Core Application Features Summary")

    add_card(slide13, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4), "Real-Time AI Response", [
        "Markdown formatting for code and tables.",
        "Context memory across chat session."
    ], "⚡")

    add_card(slide13, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4), "Automated Ticket Creation", [
        "Unique Ticket ID generation (TICK-XXXXXXXX).",
        "Direct database storage with email alerts."
    ], "🎫")

    add_card(slide13, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4), "Glassmorphism UI & Dark Mode", [
        "Modern Bootstrap 5 aesthetic.",
        "Responsive layout optimized for mobile."
    ], "💎")

    add_card(slide13, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.4), "Data Export & Analytics", [
        "Interactive Chart.js status graphs.",
        "One-click CSV ticket report download."
    ], "📈")

    add_notes(slide13, "Key features include real-time AI responses with markdown formatting and conversation memory, automated ticket creation with unique tracking IDs, modern glassmorphic dark mode styling, and built-in CSV report exporting alongside Chart.js analytics.")

    # ==========================================
    # SLIDE 14: System Advantages
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, COLOR_BG)
    add_header(slide14, "Key Advantages & Business Benefits")

    add_card(slide14, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), "Zero Friction User Access", [
        "No mandatory login screen.",
        "Instant access to support on homepage.",
        "Drastically reduces visitor bounce rates."
    ], "🚪")

    add_card(slide14, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0), "Cost & Operational Efficiency", [
        "Automates up to 80% of routine questions.",
        "Reduces human support workload.",
        "Saves infrastructure & API costs with RAG search."
    ], "💰")

    add_card(slide14, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0), "Zero-Migration Simplicity", [
        "db.create_all() auto schema creation.",
        "Works cleanly with SQLite & PostgreSQL.",
        "100% production-ready Gunicorn deployment."
    ], "⚡")

    add_notes(slide14, "The advantages of this architecture are clear: zero login friction keeps visitors engaged; automated AI and RAG FAQ search cut support costs by resolving up to 80% of repetitive tickets; and zero-migration schema auto-creation eliminates database deployment headaches.")

    # ==========================================
    # SLIDE 15: Future Enhancements
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, COLOR_BG)
    add_header(slide15, "Future System Enhancements & Scope")

    add_card(slide15, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4), "Voice & Multi-Modal Input", [
        "Integrate Speech-to-Text for real-time voice support.",
        "Support image & document attachment analysis."
    ], "🎙️")

    add_card(slide15, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4), "Multi-Language & Localization", [
        "Implement automatic language translation for international users.",
        "Localized UI themes and responses."
    ], "🌐")

    add_card(slide15, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4), "Mobile Applications", [
        "Develop native Android (Kotlin) and iOS (Swift) apps.",
        "Push notifications for ticket updates."
    ], "📱")

    add_card(slide15, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.4), "Predictive Analytics & Sentiment AI", [
        "ML model predicting potential customer churn.",
        "Automated SLA breach warning triggers."
    ], "🔮")

    add_notes(slide15, "Future enhancements include integrating voice-based interaction, multi-language localization, dedicated native mobile applications, and predictive machine learning models to forecast ticket volume and prevent SLA breaches.")

    # ==========================================
    # SLIDE 16: Conclusion
    # ==========================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16, COLOR_BG)
    add_header(slide16, "Conclusion & Summary")

    add_card(slide16, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), "Project Impact Summary", [
        "Successfully developed and deployed a modern, production-ready AI Customer Support & Ticket Management Platform.",
        "Demonstrated zero-friction guest AI assistance with RAG Knowledge Base search and real LLM integration.",
        "Implemented seamless automated ticket escalation for low-confidence queries without requiring user accounts.",
        "Delivered a secure, single-auth Admin Management Console powered by Flask-Login with live Chart.js analytics and CSV report export.",
        "Achieved 100% clean automated unit test execution and zero-migration cloud deployment on Render via Gunicorn."
    ], "🎯")

    add_notes(slide16, "In conclusion, this project demonstrates how combining Large Language Models, Knowledge Base RAG retrieval, and clean web engineering can transform customer support operations. We achieved an intuitive, zero-friction customer experience, automated escalation, and robust admin management, all built on a zero-migration, production-ready architecture.")

    # ==========================================
    # SLIDE 17: Thank You Slide
    # ==========================================
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17, COLOR_DARK_BG)

    # Thank you header box
    ty_box = slide17.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(1.8))
    tf_ty = ty_box.text_frame
    tf_ty.word_wrap = True
    
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "Thank You!"
    p_ty.font.size = Pt(44)
    p_ty.font.bold = True
    p_ty.font.color.rgb = RGBColor(255, 255, 255)
    p_ty.alignment = PP_ALIGN.CENTER
    
    p_sub = tf_ty.add_paragraph()
    p_sub.text = "Questions & Answers / Project Review Session"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_ACCENT
    p_sub.alignment = PP_ALIGN.CENTER

    # Contact Info Card
    contact_card = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(3.8), Inches(8.333), Inches(2.6))
    contact_card.fill.solid()
    contact_card.fill.fore_color.rgb = RGBColor(26, 43, 80)
    contact_card.line.color.rgb = COLOR_ACCENT

    tf_c = slide17.shapes.add_textbox(Inches(2.7), Inches(4.0), Inches(7.933), Inches(2.2)).text_frame
    tf_c.word_wrap = True

    c_info = [
        ("Presenter:", "Pamulapati Karthik (Roll No: 22N81A05K2)"),
        ("Department:", "Computer Science & Engineering (CSE)"),
        ("College:", "Sree Dattha Group of Institutions"),
        ("GitHub Repository:", "github.com/karthikpamulapti3333-debug/AI-Powered-Customer-Support-Assistant")
    ]

    for label, val in c_info:
        p = tf_c.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{label:<18} "
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_ACCENT
        
        r2 = p.add_run()
        r2.text = val
        r2.font.size = Pt(13)
        r2.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

    add_notes(slide17, "Thank you very much for your time, guidance, and attention throughout this presentation. I am now open to your questions, feedback, and project review.")

    # Save presentation
    output_filename = "AI_Customer_Support_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'")
    return output_filename

if __name__ == "__main__":
    create_deck()
